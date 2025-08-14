from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Tạo đối tượng Presentation
prs = Presentation()

# Hàm phụ để thêm slide với tiêu đề và nội dung
def add_slide(title, content, layout=1):
    slide_layout = prs.slide_layouts[layout]  # Sử dụng layout nội dung (layout 1) hoặc tiêu đề (layout 0)
    slide = prs.slides.add_slide(slide_layout)
    
    # Thêm tiêu đề
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(36)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Màu xanh đậm
    
    # Thêm nội dung
    if content:
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Xóa nội dung mặc định
        for line in content:
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.LEFT
            if line.startswith("•"):
                p.level = 1  # Định dạng dấu đầu dòng
            p.font.color.rgb = RGBColor(0, 0, 0)  # Màu đen

# Slide 1: Tiêu đề
slide_1_content = [
    "Ứng dụng Sudoku GUI sử dụng Python và Tkinter",
    "Người thực hiện: [Tên của bạn]",  # Thay bằng tên thật của bạn
    "Ngày: 22/07/2025"
]
add_slide("DỰ ÁN SUDOKU GUI", slide_1_content, layout=0)

# Slide 2: Sudoku là gì?
slide_2_content = [
    "Sudoku là trò chơi giải đố logic:",
    "• Lưới 9x9 hoặc 16x16, chia thành khối 3x3 hoặc 4x4.",
    "• Điền số (1-9) hoặc ký hiệu (A-F) vào ô trống.",
    "• Luật: Không lặp số trong hàng, cột, khối, và đường chéo (Sudoku X).",
    "Mục tiêu: Hoàn thành lưới theo luật."
]
add_slide("SUDOKU LÀ GÌ?", slide_2_content)

# Slide 3: Mục tiêu dự án
slide_3_content = [
    "Xây dựng ứng dụng Sudoku GUI:",
    "• Tạo đề Sudoku 9x9 và 16x16, hỗ trợ Sudoku X.",
    "• Hỗ trợ người chơi với gợi ý, kiểm tra, và lời giải.",
    "• Lưu/tải đề, theo dõi thời gian chơi.",
    "• Giao diện thân thiện, dễ sử dụng."
]
add_slide("MỤC TIÊU DỰ ÁN", slide_3_content)

# Slide 4: Tính năng chính
slide_4_content = [
    "Tính năng của ứng dụng:",
    "• Tạo đề ngẫu nhiên (Dễ, Trung bình, Khó).",
    "• Kiểm tra kết quả, gợi ý, hiển thị lời giải.",
    "• Lưu/tải đề từ file văn bản.",
    "• Undo/Redo để hoàn tác/làm lại.",
    "• Bộ đếm thời gian thực."
]
add_slide("TÍNH NĂNG CHÍNH", slide_4_content)

# Slide 5: Công nghệ và thuật toán
slide_5_content = [
    "Công nghệ sử dụng:",
    "• Ngôn ngữ: Python 3.",
    "• Giao diện: Thư viện Tkinter.",
    "Thuật toán:",
    "• MRV: Chọn ô có ít giá trị khả thi nhất.",
    "• Backtracking: Giải và kiểm tra lời giải.",
    "• Tạo đề: Đảm bảo 1 lời giải duy nhất."
]
add_slide("CÔNG NGHỆ & THUẬT TOÁN", slide_5_content)

# Slide 6: Kết luận
slide_6_content = [
    "Thành tựu:",
    "• Ứng dụng hoàn chỉnh, hỗ trợ 9x9, 16x16, Sudoku X.",
    "• Giao diện thân thiện, có Undo/Redo.",
    "Hạn chế:",
    "• Độ khó dựa trên số ô trống.",
    "• Chưa kiểm tra lỗi nhập liệu tức thời.",
    "Hướng phát triển: Thêm kiểm tra thời gian thực, cải thiện độ khó."
]
add_slide("KẾT LUẬN", slide_6_content)

# Lưu file PowerPoint
prs.save("Sudoku_Presentation.pptx")